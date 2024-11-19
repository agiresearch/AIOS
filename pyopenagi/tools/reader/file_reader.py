import base64
import csv
import json
import os
import zipfile
from abc import ABC, abstractmethod

import PyPDF2
import pandas as pd
import pptx
import docx

from openai import OpenAI

from pyopenagi.tools.base import BaseTool


class FileReader(BaseTool):

    def __init__(self):
        super().__init__()

    def run(self, param) -> str:
        # query = param["query"]  # Temporarily unused
        path = param["path"]

        reader = get_reader(path)
        content = reader.read(path)
        return content

    def get_tool_call_format(self):
        tool_call_format = {
            "type": "function",
            "function": {
                "name": "file_reader",
                "description": "Read the file content from the specified path.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {
                            "type": "string",
                            "description": "Relevant infromation want to retrival. If want to extract all file "
                                           "content, don't pass `query`."
                        },
                        "path": {
                            "type": "string",
                            "description": "Input file path."
                        }
                    },
                    "required": [
                        "path"
                    ]
                }
            }
        }
        return tool_call_format


READER_REGISTER = {}


def register_reader(*args):
    """
    A decorator to register a reader class into the READER_REGISTER dictionary.

    Usage:
        @register_reader(".pdf", ".docx")
        class PDFReader(Reader):
            # implement the read method
            pass

    In this example, the PDFReader class is registered to handle files with ".pdf"
    and ".docx" extensions.
    """
    def decorator(cls):
        for suffix in args:
            READER_REGISTER[suffix] = cls
        return cls

    return decorator


class Reader(ABC):

    @abstractmethod
    def read(self, path: str) -> str:
        pass


def get_reader(path: str) -> Reader:
    """
    Retrieves a reader instance based on the file extension of the given path.

    Args:
        path (str): The file path for which to get the corresponding reader.

    Returns:
        Reader: An instance of a reader class that can handle the file extension
        of the given path. This instance is expected to implement the `read` method
        to process the file content.
    """
    filename, file_extension = os.path.splitext(path)
    reader = READER_REGISTER.get(file_extension)
    return reader()


@register_reader(".jpg", ".png")
class ImageReader(Reader):
    """
    Reads image files and recognizes their content using the OpenAI API.
    """

    def __init__(self):
        self.client = OpenAI()

    def read(self, path: str) -> str:
        with open(path, "rb") as image_file:
            base64_img = base64.b64encode(image_file.read()).decode("utf-8")

        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "What's in this Image?"},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_img}",
                            },
                        },
                    ],
                }
            ],
            max_tokens=300,
        )
        return response.choices[0].message.content


@register_reader(".pdf")
class PDFReader(Reader):
    """
    Reads PDF files and returns the content of each page as plain text.
    """

    def read(self, path: str) -> str:
        content = ""
        pdf = PyPDF2.PdfReader(path)
        for index, page in enumerate(pdf.pages):
            content += f"Page {index + 1}:\n" + page.extract_text()
        return content


@register_reader(".pptx")
class PPTReader(Reader):
    """
    Reads PowerPoint files (.pptx) and returns the content of each slide as plain text.

    The content of each slide is concatenated and returned as a single string.

    Note: This reader does not support PPT files that contain images or other
    non-text content.
    """

    def read(self, path: str) -> str:
        content = ""
        ppt = pptx.Presentation(path)
        for index, slide in enumerate(ppt.slides):
            content += f"Slide {index + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text
        return content


@register_reader(".txt", ".pdb")
class TextReader(Reader):
    """
    Reads plain text files (.txt and .pdb) and returns the content as a string.

    The content of the file is read and returned as a single string.
    """

    def read(self, path: str) -> str:
        content = ""
        with open(path, "r", encoding="utf-8") as file:
            content += file.read()
        return content


@register_reader(".mp3")
class AudioReader(Reader):
    """
    Reads audio files (.mp3) and transcribes the content using the Whisper API.

    This reader uses the Whisper API to transcribe the audio file and returns the
    transcription as a string.

    Note: This reader requires an OpenAI API key to be set in the environment.
    """

    def __init__(self):
        self.client = OpenAI()

    def read(self, path: str) -> str:
        content = ""
        with open(path, "rb") as audio_file:
            response = self.client.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file
            )
            content += response.text

        return content


@register_reader(".xlsx")
class ExcelReader(Reader):
    """
    Reads Excel files (.xlsx) and returns the content as a string.

    The content of the file is read into a pandas DataFrame and then converted
    to a string using the `to_string` method.

    Note: This reader requires the `pandas` library to be installed.
    """

    def read(self, path: str) -> str:
        content = ""
        data = pd.read_excel(path)
        content += data.to_string()
        return content


@register_reader(".json", "jsonl")
class JsonReader(Reader):
    """
    Reads JSON and JSONL files and returns their content as a string.

    For JSON files, the content is loaded as a dictionary and converted to a string.
    For JSONL files, each line is loaded as a JSON object and the list of these
    objects is converted to a string.
    """

    def read(self, path):
        content = ""
        with open(path, "r", encoding="utf-8") as file:
            if path.endswith(".json"):
                content += str(json.load(file))
            elif path.endswith(".jsonl"):
                content += str([json.loads(line) for line in file])
        return content


@register_reader(".docx")
class DocxReader(Reader):
    """
    Reads DOCX files and returns the content as plain text.

    This reader processes each paragraph in the DOCX file, appending
    it to the content string with a page indicator.
    """

    def read(self, path: str) -> str:
        content = ""
        doc = docx.Document(path)
        for index, page in enumerate(doc.paragraphs):
            content += f"Page {index + 1}:\n" + page.text + "\n"
        return content


@register_reader(".py")
class PythonReader(Reader):
    """
    Reads Python files (.py) and returns their content as a string.

    The content of the file is read and returned as a single string.
    """

    def read(self, path: str) -> str:
        content = ""
        with open(path, "r", encoding="utf-8") as file:
            content += file.read()
        return content


@register_reader(".zip")
class ZipReader(Reader):
    """
    Reads the content of a ZIP file.

    This reader will extract the ZIP file to the same directory and then
    read the content of each file in the ZIP file. The content of each file
    is processed by the corresponding reader and the results are concatenated
    together with each file's name and content separated by a blank line.
    """

    def read(self, path: str) -> str:
        content = ""
        with zipfile.ZipFile(path, "r") as zip_file:
            extract_dir = path[:-4] + '/'
            zip_file.extractall(extract_dir)
            for file_name in zip_file.namelist():
                content += f"File {file_name}:\n"

                sub_reader = get_reader(extract_dir + file_name)
                content += sub_reader.read(extract_dir + file_name) + "\n"

        return content


@register_reader(".csv")
class CSVReader(Reader):
    """
    Reads CSV files and returns their content as a string.

    The content of the CSV file is read into a list of rows, where each row is a list of values.
    The list is then converted to a string and returned.
    """

    def read(self, path: str) -> str:
        content = ""
        with open(path, newline='', encoding="utf-8") as file:
            reader = csv.reader(file)
            data = [row for row in reader]
            content += str(data)
        return content


if __name__ == "__main__":
    param = {
        "path": "<absolute path>"
    }
    reader = FileReader()
    content = reader.run(param)
    print(f"File Path: {param['path']}\n"
          f"<Content>\n "
          f"{content}\n "
          f"</Content>")
